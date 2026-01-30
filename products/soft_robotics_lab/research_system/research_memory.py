#!/usr/bin/env python3
"""
OMEGA Research Memory
=====================

Private research OS for a PI's entire corpus:
- Ingest PDFs, notes, slides, protocols
- Structured retrieval (keyword + semantic)
- Research objects (hypothesis, experiment, result cards)
- Insight engines (contradictions, discriminators, fragility)
- Weekly briefs with full traceability

Every insight is anchored to source snippets or flagged as speculation.

Usage:
    from research_memory import ResearchMemory
    memory = ResearchMemory("jon_lab")
    memory.ingest("papers/")
    memory.query_with_sources("What do we know about variable stiffness grippers?")
"""

import sys
import json
import sqlite3
import hashlib
from pathlib import Path
from datetime import datetime
from dataclasses import dataclass, field, asdict
from typing import List, Dict, Optional, Tuple
import re

_repo_root = Path(__file__).resolve().parent.parent.parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))


# ═══════════════════════════════════════════════════════════════════════════
# DATA MODELS
# ═══════════════════════════════════════════════════════════════════════════

@dataclass
class DocumentChunk:
    """A chunk of a document with full provenance."""
    id: str
    doc_id: str
    doc_path: str
    doc_type: str  # paper, note, protocol, slide, grant

    # Content
    text: str
    section: str  # e.g., "Methods", "Results", "Figure 3 caption"
    page: Optional[int] = None

    # Metadata
    project: str = ""
    authors: List[str] = field(default_factory=list)
    date: str = ""

    # Entities extracted
    entities: List[str] = field(default_factory=list)  # materials, methods, organisms

    # Embedding (stored separately)
    embedding_id: Optional[str] = None


@dataclass
class HypothesisCard:
    """A structured hypothesis with evidence links."""
    id: str

    # Core claim
    claim: str  # Mechanistic, falsifiable
    scope: str  # What conditions this applies to
    assumptions: List[str]

    # Competition
    competing_hypotheses: List[str]  # IDs of competing cards

    # Predictions
    predictions: List[str]  # What we'd observe if true

    # Evidence
    supporting_evidence: List[str]  # Chunk IDs
    weakening_evidence: List[str]  # Chunk IDs

    # Status
    status: str  # active, weakened, falsified, merged
    confidence: float  # 0-1

    # Next steps
    next_discriminator: str  # What experiment would test this

    # Metadata
    created: str = ""
    updated: str = ""
    created_by: str = ""

    def to_dict(self) -> Dict:
        return asdict(self)


@dataclass
class ExperimentCard:
    """A structured experiment record."""
    id: str

    # Purpose
    aim: str
    hypotheses_tested: List[str]  # Hypothesis IDs this discriminates

    # Design
    design_summary: str
    key_variables: List[str]
    controls: List[str]

    # Predictions
    expected_outcomes: Dict[str, str]  # hypothesis_id -> expected outcome

    # Results
    result_ids: List[str]  # Link to ResultCards

    # Status
    status: str  # planned, running, completed, failed
    failure_modes: List[str]

    # Metadata
    created: str = ""
    updated: str = ""

    def to_dict(self) -> Dict:
        return asdict(self)


@dataclass
class ResultCard:
    """A structured result with uncertainty."""
    id: str
    experiment_id: str

    # Summary
    summary: str

    # Quality
    uncertainty: str  # low, medium, high
    quality_flags: List[str]  # weak_control, small_n, confound, etc.

    # Impact
    hypotheses_updated: List[str]  # Which hypothesis edges this updates
    direction: str  # supports, weakens, neutral

    # Replication
    replicated: bool = False
    replication_notes: str = ""

    # Negative result (first-class)
    is_negative: bool = False
    negative_insight: str = ""  # What we learned from null result

    # Evidence
    source_chunks: List[str] = field(default_factory=list)  # Chunk IDs

    # Metadata
    created: str = ""

    def to_dict(self) -> Dict:
        return asdict(self)


@dataclass
class InsightBrief:
    """Weekly insight brief for the lab."""
    id: str
    week: str
    generated: str

    # Sections
    what_changed: List[str]
    top_contradictions: List[str]
    top_discriminators: List[str]
    fragility_alerts: List[str]
    portfolio_shifts: List[str]
    questions_for_pi: List[str]

    # Audit
    sources_used: List[str]  # Chunk IDs
    speculation_flags: List[str]  # Which insights are speculative

    def to_markdown(self) -> str:
        md = f"""# Weekly Insight Brief

**Week:** {self.week}
**Generated:** {self.generated}

---

## 1. What Changed

"""
        for item in self.what_changed:
            md += f"- {item}\n"

        md += """
---

## 2. Top Contradictions / Tensions

"""
        for item in self.top_contradictions:
            md += f"- {item}\n"

        md += """
---

## 3. Top Discriminators (Smallest Tests)

"""
        for item in self.top_discriminators:
            md += f"- {item}\n"

        md += """
---

## 4. Fragility Alerts

"""
        for item in self.fragility_alerts:
            md += f"- {item}\n"

        md += """
---

## 5. Portfolio Shifts

"""
        for item in self.portfolio_shifts:
            md += f"- {item}\n"

        md += """
---

## 6. Questions for PI

"""
        for i, q in enumerate(self.questions_for_pi, 1):
            md += f"{i}. {q}\n"

        if self.speculation_flags:
            md += """
---

## Speculation Flags

*The following insights are flagged as speculative (not fully grounded in sources):*

"""
            for flag in self.speculation_flags:
                md += f"- {flag}\n"

        md += f"""
---

*Brief ID: {self.id} | Sources: {len(self.sources_used)} chunks*
"""
        return md


# ═══════════════════════════════════════════════════════════════════════════
# RESEARCH MEMORY SYSTEM
# ═══════════════════════════════════════════════════════════════════════════

class ResearchMemory:
    """
    Private research memory for a lab.

    Combines:
    - Document ingestion + chunking
    - Keyword + semantic search
    - Structured research objects
    - Insight engines
    """

    def __init__(self, lab_name: str, base_path: Path = None):
        self.lab_name = lab_name
        self.base_path = base_path or Path(__file__).resolve().parent / "data" / lab_name
        self.base_path.mkdir(parents=True, exist_ok=True)

        self.db_path = self.base_path / "research_memory.db"
        self._init_db()

        self.embeddings: Dict[str, List[float]] = {}
        self._embedding_model = None  # Lazy-loaded SentenceTransformer or False if unavailable

    def _init_db(self) -> None:
        """Initialize database schema."""
        conn = sqlite3.connect(self.db_path)
        conn.executescript("""
            CREATE TABLE IF NOT EXISTS documents (
                id TEXT PRIMARY KEY,
                path TEXT,
                doc_type TEXT,
                project TEXT,
                authors TEXT,
                date TEXT,
                hash TEXT,
                ingested_at TEXT
            );

            CREATE TABLE IF NOT EXISTS chunks (
                id TEXT PRIMARY KEY,
                doc_id TEXT,
                text TEXT,
                section TEXT,
                page INTEGER,
                entities TEXT,
                FOREIGN KEY (doc_id) REFERENCES documents(id)
            );

            CREATE TABLE IF NOT EXISTS hypotheses (
                id TEXT PRIMARY KEY,
                data TEXT,
                status TEXT,
                confidence REAL,
                updated TEXT
            );

            CREATE TABLE IF NOT EXISTS experiments (
                id TEXT PRIMARY KEY,
                data TEXT,
                status TEXT,
                updated TEXT
            );

            CREATE TABLE IF NOT EXISTS results (
                id TEXT PRIMARY KEY,
                experiment_id TEXT,
                data TEXT,
                is_negative INTEGER,
                created TEXT
            );

            CREATE TABLE IF NOT EXISTS briefs (
                id TEXT PRIMARY KEY,
                week TEXT,
                data TEXT,
                generated TEXT
            );

            CREATE TABLE IF NOT EXISTS audit_log (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                timestamp TEXT,
                action TEXT,
                object_type TEXT,
                object_id TEXT,
                user TEXT,
                details TEXT
            );

            CREATE TABLE IF NOT EXISTS chunk_embeddings (
                chunk_id TEXT PRIMARY KEY,
                embedding TEXT,
                model_name TEXT,
                FOREIGN KEY (chunk_id) REFERENCES chunks(id)
            );

            CREATE INDEX IF NOT EXISTS idx_chunks_doc ON chunks(doc_id);
            CREATE INDEX IF NOT EXISTS idx_chunk_embeddings_model ON chunk_embeddings(model_name);
            CREATE INDEX IF NOT EXISTS idx_hypotheses_status ON hypotheses(status);
            CREATE INDEX IF NOT EXISTS idx_results_experiment ON results(experiment_id);
        """)
        conn.commit()
        conn.close()

    # ═══════════════════════════════════════════════════════════════════════
    # INGESTION
    # ═══════════════════════════════════════════════════════════════════════

    def ingest(self, path: str, project: str = "default", doc_type: str = None) -> int:
        """
        Ingest documents from a path.

        Returns number of chunks created.
        """
        path = Path(path)
        total_chunks = 0

        if path.is_file():
            total_chunks += self._ingest_file(path, project, doc_type)
        elif path.is_dir():
            for file_path in path.rglob("*"):
                if file_path.is_file() and file_path.suffix.lower() in [
                    ".pdf", ".md", ".txt", ".pptx"
                ]:
                    total_chunks += self._ingest_file(file_path, project, doc_type)

        return total_chunks

    def _ingest_file(
        self, file_path: Path, project: str, doc_type: str = None
    ) -> int:
        """Ingest a single file."""
        if doc_type is None:
            doc_type = self._detect_doc_type(file_path)

        file_hash = self._hash_file(file_path)

        conn = sqlite3.connect(self.db_path)
        cursor = conn.execute("SELECT id FROM documents WHERE hash = ?", (file_hash,))
        if cursor.fetchone():
            conn.close()
            return 0

        doc_id = f"doc_{datetime.now().strftime('%Y%m%d%H%M%S')}_{file_hash[:8]}"

        text = self._extract_text(file_path)
        chunks = self._chunk_text(text, doc_id, str(file_path), doc_type)

        conn.execute(
            """
            INSERT INTO documents (id, path, doc_type, project, authors, date, hash, ingested_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                doc_id,
                str(file_path),
                doc_type,
                project,
                "",
                "",
                file_hash,
                datetime.now().isoformat(),
            ),
        )

        for chunk in chunks:
            conn.execute(
                """
                INSERT INTO chunks (id, doc_id, text, section, page, entities)
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                (
                    chunk.id,
                    chunk.doc_id,
                    chunk.text,
                    chunk.section,
                    chunk.page,
                    json.dumps(chunk.entities),
                ),
            )

        conn.commit()
        conn.close()

        self._store_embeddings_for_chunks(chunks)

        self._audit_log(
            "ingest", "document", doc_id, f"Ingested {file_path.name}, {len(chunks)} chunks"
        )

        return len(chunks)

    def _detect_doc_type(self, path: Path) -> str:
        """Detect document type from filename/content."""
        name = path.name.lower()

        if "protocol" in name:
            return "protocol"
        elif "note" in name or "meeting" in name:
            return "note"
        elif "grant" in name or "proposal" in name:
            return "grant"
        elif path.suffix == ".pptx":
            return "slide"
        elif path.suffix == ".pdf":
            return "paper"
        else:
            return "note"

    def _hash_file(self, path: Path) -> str:
        """Create hash of file content."""
        return hashlib.md5(path.read_bytes()).hexdigest()

    def _extract_text(self, path: Path) -> str:
        """Extract text from file."""
        suffix = path.suffix.lower()

        if suffix in (".txt", ".md"):
            return path.read_text(encoding="utf-8", errors="ignore")

        elif suffix == ".pdf":
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(path)
                text = ""
                for page in doc:
                    text += page.get_text() + "\n\n"
                return text
            except ImportError:
                return f"[PDF: {path.name} - text extraction requires PyMuPDF]"

        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    text += "\n---\n"
                return text
            except ImportError:
                return f"[PPTX: {path.name} - extraction requires python-pptx]"

        return ""

    def _chunk_text(
        self,
        text: str,
        doc_id: str,
        doc_path: str,
        doc_type: str,
    ) -> List[DocumentChunk]:
        """Chunk text by structure (sections, paragraphs)."""
        chunks = []

        sections = re.split(r"\n(?=#{1,3}\s|\n\n)", text)

        for i, section_text in enumerate(sections):
            if not section_text.strip():
                continue

            lines = section_text.strip().split("\n")
            if lines[0].startswith("#"):
                section_title = lines[0].lstrip("#").strip()
                content = "\n".join(lines[1:])
            else:
                section_title = f"Section {i+1}"
                content = section_text

            if len(content) > 1000:
                paragraphs = content.split("\n\n")
                for j, para in enumerate(paragraphs):
                    if para.strip():
                        chunk_id = f"{doc_id}_chunk_{i}_{j}"
                        chunks.append(
                            DocumentChunk(
                                id=chunk_id,
                                doc_id=doc_id,
                                doc_path=doc_path,
                                doc_type=doc_type,
                                text=para.strip()[:2000],
                                section=section_title,
                                page=None,
                                entities=self._extract_entities(para),
                            )
                        )
            else:
                chunk_id = f"{doc_id}_chunk_{i}"
                chunks.append(
                    DocumentChunk(
                        id=chunk_id,
                        doc_id=doc_id,
                        doc_path=doc_path,
                        doc_type=doc_type,
                        text=content.strip()[:2000],
                        section=section_title,
                        page=None,
                        entities=self._extract_entities(content),
                    )
                )

        return chunks

    def _extract_entities(self, text: str) -> List[str]:
        """Extract key entities (materials, methods, etc.)."""
        entities = []

        patterns = [
            r"\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+\b",
            r"\b\d+\s*(?:mm|cm|m|kg|g|N|Pa|Hz|kHz|MHz)\b",
            r"\bPLA|ABS|PDMS|silicone|rubber|steel\b",
        ]

        for pattern in patterns:
            matches = re.findall(pattern, text)
            entities.extend(matches[:5])

        return list(set(entities))[:10]

    # ═══════════════════════════════════════════════════════════════════════
    # EMBEDDINGS & SEMANTIC SEARCH
    # ═══════════════════════════════════════════════════════════════════════

    EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2"

    def _get_embedding_model(self):
        """Lazy-load SentenceTransformer; return None if unavailable."""
        if self._embedding_model is None:
            try:
                from sentence_transformers import SentenceTransformer
                self._embedding_model = SentenceTransformer(self.EMBEDDING_MODEL_NAME)
            except ImportError:
                self._embedding_model = False
        return self._embedding_model if self._embedding_model else None

    def _embed_text(self, text: str) -> Optional[List[float]]:
        """Embed text; return None if model unavailable."""
        model = self._get_embedding_model()
        if model is None:
            return None
        emb = model.encode(text, convert_to_numpy=True, normalize_embeddings=True)
        return emb.tolist()

    def _store_embeddings_for_chunks(self, chunks: List[DocumentChunk]) -> None:
        """Compute and store embeddings for ingested chunks."""
        model = self._get_embedding_model()
        if model is None:
            return
        conn = sqlite3.connect(self.db_path)
        for chunk in chunks:
            emb = self._embed_text(chunk.text)
            if emb is not None:
                conn.execute(
                    """
                    INSERT OR REPLACE INTO chunk_embeddings (chunk_id, embedding, model_name)
                    VALUES (?, ?, ?)
                    """,
                    (chunk.id, json.dumps(emb), self.EMBEDDING_MODEL_NAME),
                )
        conn.commit()
        conn.close()

    def reindex_embeddings(self) -> int:
        """
        (Re)compute embeddings for all chunks. Use after installing sentence-transformers.
        Returns number of chunks indexed.
        """
        if self._get_embedding_model() is None:
            return 0
        conn = sqlite3.connect(self.db_path)
        cursor = conn.execute(
            """
            SELECT c.id, c.doc_id, d.path, d.doc_type, c.text, c.section, c.page, c.entities
            FROM chunks c
            JOIN documents d ON c.doc_id = d.id
            """
        )
        chunks = []
        for row in cursor.fetchall():
            chunks.append(
                DocumentChunk(
                    id=row[0],
                    doc_id=row[1],
                    doc_path=row[2],
                    doc_type=row[3],
                    text=row[4],
                    section=row[5],
                    page=row[6],
                    entities=json.loads(row[7]) if row[7] else [],
                )
            )
        conn.close()
        self._store_embeddings_for_chunks(chunks)
        return len(chunks)

    def _search_semantic(
        self,
        query: str,
        project: str = None,
        doc_type: str = None,
        limit: int = 10,
    ) -> List[DocumentChunk]:
        """Semantic search using stored embeddings. Returns [] if unavailable."""
        model = self._get_embedding_model()
        if model is None:
            return []
        q_emb = self._embed_text(query)
        if q_emb is None:
            return []
        conn = sqlite3.connect(self.db_path)
        cursor = conn.execute(
            "SELECT chunk_id, embedding FROM chunk_embeddings WHERE model_name = ?",
            (self.EMBEDDING_MODEL_NAME,),
        )
        rows = cursor.fetchall()
        conn.close()
        if not rows:
            return []
        q_vec = q_emb
        scored = []
        for chunk_id, emb_json in rows:
            c_vec = json.loads(emb_json)
            sim = sum(a * b for a, b in zip(q_vec, c_vec))
            scored.append((chunk_id, sim))
        scored.sort(key=lambda x: -x[1])
        top_ids = [cid for cid, _ in scored[: limit * 2]]
        if not top_ids:
            return []
        conn = sqlite3.connect(self.db_path)
        placeholders = ",".join("?" * len(top_ids))
        sql = f"""
            SELECT c.id, c.doc_id, d.path, d.doc_type, c.text, c.section, c.page, c.entities
            FROM chunks c
            JOIN documents d ON c.doc_id = d.id
            WHERE c.id IN ({placeholders})
            """
        params = list(top_ids)
        if project:
            sql += " AND d.project = ?"
            params.append(project)
        if doc_type:
            sql += " AND d.doc_type = ?"
            params.append(doc_type)
        cursor = conn.execute(sql, params)
        rows_by_id = {row[0]: row for row in cursor.fetchall()}
        conn.close()
        order = [cid for cid in top_ids if cid in rows_by_id][:limit]
        chunks = []
        for cid in order:
            row = rows_by_id[cid]
            chunks.append(
                DocumentChunk(
                    id=row[0],
                    doc_id=row[1],
                    doc_path=row[2],
                    doc_type=row[3],
                    text=row[4],
                    section=row[5],
                    page=row[6],
                    entities=json.loads(row[7]) if row[7] else [],
                )
            )
        return chunks

    # ═══════════════════════════════════════════════════════════════════════
    # SEARCH
    # ═══════════════════════════════════════════════════════════════════════

    def search(
        self,
        query: str,
        project: str = None,
        doc_type: str = None,
        limit: int = 10,
        use_semantic: bool = True,
    ) -> List[DocumentChunk]:
        """
        Search for chunks matching query.

        Uses semantic search when embeddings are available (sentence-transformers);
        otherwise falls back to keyword (LIKE) search.
        """
        if use_semantic:
            chunks = self._search_semantic(
                query, project=project, doc_type=doc_type, limit=limit
            )
            if chunks:
                return chunks
        conn = sqlite3.connect(self.db_path)
        sql = """
            SELECT c.id, c.doc_id, d.path, d.doc_type, c.text, c.section, c.page, c.entities
            FROM chunks c
            JOIN documents d ON c.doc_id = d.id
            WHERE c.text LIKE ?
            """
        params = [f"%{query}%"]

        if project:
            sql += " AND d.project = ?"
            params.append(project)

        if doc_type:
            sql += " AND d.doc_type = ?"
            params.append(doc_type)

        sql += f" LIMIT {limit}"

        cursor = conn.execute(sql, params)

        chunks = []
        for row in cursor.fetchall():
            chunks.append(
                DocumentChunk(
                    id=row[0],
                    doc_id=row[1],
                    doc_path=row[2],
                    doc_type=row[3],
                    text=row[4],
                    section=row[5],
                    page=row[6],
                    entities=json.loads(row[7]) if row[7] else [],
                )
            )

        conn.close()
        return chunks

    def query_with_sources(
        self, question: str, project: str = None, use_semantic: bool = True
    ) -> Dict:
        """
        Answer a question with full source citations.

        Returns answer + source chunks + speculation flags.
        """
        chunks = self.search(
            question, project=project, limit=5, use_semantic=use_semantic
        )

        if not chunks:
            return {
                "answer": "No relevant sources found in the research memory.",
                "sources": [],
                "speculation": True,
                "speculation_reason": "No source material available",
            }

        answer = f"Based on {len(chunks)} sources:\n\n"
        for c in chunks:
            answer += f"• From {Path(c.doc_path).name} ({c.section}): {c.text[:200]}...\n\n"

        return {
            "answer": answer,
            "sources": [
                {
                    "id": c.id,
                    "doc": c.doc_path,
                    "section": c.section,
                    "text": c.text[:500],
                }
                for c in chunks
            ],
            "speculation": False,
            "speculation_reason": None,
        }

    # ═══════════════════════════════════════════════════════════════════════
    # RESEARCH OBJECTS
    # ═══════════════════════════════════════════════════════════════════════

    def create_hypothesis(
        self,
        claim: str,
        scope: str,
        assumptions: List[str],
        predictions: List[str],
        created_by: str = "unknown",
    ) -> HypothesisCard:
        """Create a new hypothesis card."""
        hyp_id = f"hyp_{datetime.now().strftime('%Y%m%d%H%M%S')}"

        card = HypothesisCard(
            id=hyp_id,
            claim=claim,
            scope=scope,
            assumptions=assumptions,
            competing_hypotheses=[],
            predictions=predictions,
            supporting_evidence=[],
            weakening_evidence=[],
            status="active",
            confidence=0.5,
            next_discriminator="",
            created=datetime.now().isoformat(),
            updated=datetime.now().isoformat(),
            created_by=created_by,
        )

        conn = sqlite3.connect(self.db_path)
        conn.execute(
            """
            INSERT INTO hypotheses (id, data, status, confidence, updated)
            VALUES (?, ?, ?, ?, ?)
            """,
            (
                card.id,
                json.dumps(card.to_dict()),
                card.status,
                card.confidence,
                card.updated,
            ),
        )
        conn.commit()
        conn.close()

        self._audit_log("create", "hypothesis", card.id, f"Created: {claim[:50]}")

        return card

    def get_hypotheses(self, status: str = None) -> List[HypothesisCard]:
        """Get all hypotheses, optionally filtered by status."""
        conn = sqlite3.connect(self.db_path)

        if status:
            cursor = conn.execute(
                "SELECT data FROM hypotheses WHERE status = ?", (status,)
            )
        else:
            cursor = conn.execute("SELECT data FROM hypotheses")

        cards = []
        for row in cursor.fetchall():
            data = json.loads(row[0])
            cards.append(HypothesisCard(**data))

        conn.close()
        return cards

    def update_hypothesis(self, hyp_id: str, **updates) -> bool:
        """Update a hypothesis card."""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.execute("SELECT data FROM hypotheses WHERE id = ?", (hyp_id,))
        row = cursor.fetchone()

        if not row:
            conn.close()
            return False

        data = json.loads(row[0])
        data.update(updates)
        data["updated"] = datetime.now().isoformat()

        conn.execute(
            """
            UPDATE hypotheses SET data = ?, status = ?, confidence = ?, updated = ?
            WHERE id = ?
            """,
            (
                json.dumps(data),
                data.get("status", "active"),
                data.get("confidence", 0.5),
                data["updated"],
                hyp_id,
            ),
        )
        conn.commit()
        conn.close()

        self._audit_log(
            "update", "hypothesis", hyp_id, f"Updated: {list(updates.keys())}"
        )

        return True

    def add_evidence(self, hyp_id: str, chunk_id: str, direction: str) -> bool:
        """Add evidence link to a hypothesis."""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.execute("SELECT data FROM hypotheses WHERE id = ?", (hyp_id,))
        row = cursor.fetchone()

        if not row:
            conn.close()
            return False

        data = json.loads(row[0])

        if direction == "supports":
            if chunk_id not in data["supporting_evidence"]:
                data["supporting_evidence"].append(chunk_id)
        elif direction == "weakens":
            if chunk_id not in data["weakening_evidence"]:
                data["weakening_evidence"].append(chunk_id)

        data["updated"] = datetime.now().isoformat()

        conn.execute(
            "UPDATE hypotheses SET data = ?, updated = ? WHERE id = ?",
            (json.dumps(data), data["updated"], hyp_id),
        )
        conn.commit()
        conn.close()

        return True

    # ═══════════════════════════════════════════════════════════════════════
    # INSIGHT ENGINES
    # ═══════════════════════════════════════════════════════════════════════

    def find_contradictions(self) -> List[Dict]:
        """Find contradictions between hypotheses or claims."""
        hypotheses = self.get_hypotheses(status="active")
        contradictions = []

        for i, h1 in enumerate(hypotheses):
            for h2 in hypotheses[i + 1 :]:
                overlap = any(
                    word in h2.scope.lower() for word in h1.scope.lower().split()
                )

                if overlap:
                    for p1 in h1.predictions:
                        for p2 in h2.predictions:
                            if self._predictions_conflict(p1, p2):
                                contradictions.append({
                                    "hypothesis_1": h1.id,
                                    "claim_1": h1.claim,
                                    "hypothesis_2": h2.id,
                                    "claim_2": h2.claim,
                                    "conflict": f"Predictions '{p1}' vs '{p2}'",
                                    "suggested_discriminator": (
                                        f"Test that distinguishes {p1} from {p2}"
                                    ),
                                })

        return contradictions

    def _predictions_conflict(self, p1: str, p2: str) -> bool:
        """Check if two predictions conflict (simplified)."""
        opposites = [
            ("increase", "decrease"),
            ("higher", "lower"),
            ("more", "less"),
            ("positive", "negative"),
        ]

        p1_lower = p1.lower()
        p2_lower = p2.lower()

        for word1, word2 in opposites:
            if (word1 in p1_lower and word2 in p2_lower) or (
                word2 in p1_lower and word1 in p2_lower
            ):
                return True

        return False

    def surface_assumptions(self, hyp_id: str) -> List[Dict]:
        """Surface implicit assumptions in a hypothesis."""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.execute("SELECT data FROM hypotheses WHERE id = ?", (hyp_id,))
        row = cursor.fetchone()

        if not row:
            conn.close()
            return []

        data = json.loads(row[0])
        conn.close()

        assumptions = []

        for a in data.get("assumptions", []):
            assumptions.append({
                "assumption": a,
                "type": "explicit",
                "testable": True,
                "risk": "medium",
            })

        claim = data.get("claim", "")

        implicit_patterns = [
            ("will", "Assumes future conditions remain similar to present"),
            ("always", "Assumes no exceptions exist"),
            ("causes", "Assumes causation, not just correlation"),
            ("because", "Assumes mechanism is correctly identified"),
        ]

        for pattern, assumption in implicit_patterns:
            if pattern in claim.lower():
                assumptions.append({
                    "assumption": assumption,
                    "type": "implicit",
                    "testable": True,
                    "risk": "high",
                })

        return assumptions

    def generate_discriminators(self, hyp_ids: List[str]) -> List[Dict]:
        """Generate experiments that discriminate between hypotheses."""
        hypotheses = []
        conn = sqlite3.connect(self.db_path)

        for hyp_id in hyp_ids:
            cursor = conn.execute("SELECT data FROM hypotheses WHERE id = ?", (hyp_id,))
            row = cursor.fetchone()
            if row:
                hypotheses.append(json.loads(row[0]))

        conn.close()

        if len(hypotheses) < 2:
            return []

        discriminators = []

        all_predictions = {}
        for h in hypotheses:
            for p in h.get("predictions", []):
                if p not in all_predictions:
                    all_predictions[p] = []
                all_predictions[p].append(h["id"])

        for pred, hyp_list in all_predictions.items():
            if len(hyp_list) == 1:
                discriminators.append({
                    "experiment": f"Test: {pred}",
                    "discriminates": hyp_list[0],
                    "if_true": f"Supports {hyp_list[0]}",
                    "if_false": f"Weakens {hyp_list[0]}",
                    "cost": "unknown",
                    "information_gain": "high",
                })

        return discriminators

    def generate_weekly_brief(self) -> InsightBrief:
        """Generate the weekly insight brief."""
        brief_id = f"brief_{datetime.now().strftime('%Y%m%d')}"
        week = datetime.now().strftime("%Y-W%W")

        hypotheses = self.get_hypotheses()
        contradictions = self.find_contradictions()

        what_changed = []
        for h in hypotheses:
            if h.updated:
                try:
                    updated = datetime.fromisoformat(h.updated)
                    if (datetime.now() - updated).days <= 7:
                        what_changed.append(
                            f"Hypothesis '{h.claim[:50]}' updated - status: {h.status}"
                        )
                except (ValueError, TypeError):
                    pass

        top_contradictions = [
            f"{c['claim_1'][:30]} vs {c['claim_2'][:30]}: {c['conflict']}"
            for c in contradictions[:3]
        ]

        active_hyps = [h.id for h in hypotheses if h.status == "active"]
        discriminators = self.generate_discriminators(active_hyps[:4])
        top_discriminators = [d["experiment"] for d in discriminators[:3]]

        fragility_alerts = []
        for h in hypotheses:
            if h.status == "active" and len(h.supporting_evidence) < 2:
                fragility_alerts.append(
                    f"'{h.claim[:40]}' has thin evidence ({len(h.supporting_evidence)} sources)"
                )

        portfolio_shifts = []
        for h in hypotheses:
            if h.status == "weakened":
                portfolio_shifts.append(f"Consider killing: '{h.claim[:40]}'")
            elif h.confidence > 0.8 and h.status == "active":
                portfolio_shifts.append(f"High confidence, pursue: '{h.claim[:40]}'")

        questions = [
            "Which hypotheses should we prioritize this week?",
            "Any new competing hypotheses from recent literature?",
            "Resource allocation: which experiments are highest priority?",
        ]

        brief = InsightBrief(
            id=brief_id,
            week=week,
            generated=datetime.now().isoformat(),
            what_changed=what_changed or ["No significant changes this week"],
            top_contradictions=top_contradictions or ["No contradictions detected"],
            top_discriminators=top_discriminators or ["No clear discriminators identified"],
            fragility_alerts=fragility_alerts or ["No fragility alerts"],
            portfolio_shifts=portfolio_shifts or ["Portfolio stable"],
            questions_for_pi=questions,
            sources_used=[],
            speculation_flags=["Weekly analysis is partially automated and should be reviewed"],
        )

        conn = sqlite3.connect(self.db_path)
        conn.execute(
            """
            INSERT OR REPLACE INTO briefs (id, week, data, generated)
            VALUES (?, ?, ?, ?)
            """,
            (brief.id, brief.week, json.dumps(asdict(brief)), brief.generated),
        )
        conn.commit()
        conn.close()

        return brief

    # ═══════════════════════════════════════════════════════════════════════
    # AUDIT
    # ═══════════════════════════════════════════════════════════════════════

    def _audit_log(
        self,
        action: str,
        object_type: str,
        object_id: str,
        details: str,
        user: str = "system",
    ) -> None:
        """Log an action to the audit trail."""
        conn = sqlite3.connect(self.db_path)
        conn.execute(
            """
            INSERT INTO audit_log (timestamp, action, object_type, object_id, user, details)
            VALUES (?, ?, ?, ?, ?, ?)
            """,
            (datetime.now().isoformat(), action, object_type, object_id, user, details),
        )
        conn.commit()
        conn.close()

    def get_audit_log(
        self, object_id: str = None, limit: int = 50
    ) -> List[Dict]:
        """Get audit log entries."""
        conn = sqlite3.connect(self.db_path)

        if object_id:
            cursor = conn.execute(
                """
                SELECT timestamp, action, object_type, object_id, user, details
                FROM audit_log WHERE object_id = ? ORDER BY timestamp DESC LIMIT ?
                """,
                (object_id, limit),
            )
        else:
            cursor = conn.execute(
                """
                SELECT timestamp, action, object_type, object_id, user, details
                FROM audit_log ORDER BY timestamp DESC LIMIT ?
                """,
                (limit,),
            )

        logs = []
        for row in cursor.fetchall():
            logs.append({
                "timestamp": row[0],
                "action": row[1],
                "object_type": row[2],
                "object_id": row[3],
                "user": row[4],
                "details": row[5],
            })

        conn.close()
        return logs

    def get_stats(self) -> Dict:
        """Get memory statistics."""
        conn = sqlite3.connect(self.db_path)

        stats = {}

        cursor = conn.execute("SELECT COUNT(*) FROM documents")
        stats["documents"] = cursor.fetchone()[0]

        cursor = conn.execute("SELECT COUNT(*) FROM chunks")
        stats["chunks"] = cursor.fetchone()[0]

        cursor = conn.execute("SELECT COUNT(*) FROM hypotheses")
        stats["hypotheses"] = cursor.fetchone()[0]

        cursor = conn.execute(
            "SELECT COUNT(*) FROM hypotheses WHERE status = 'active'"
        )
        stats["active_hypotheses"] = cursor.fetchone()[0]

        cursor = conn.execute("SELECT COUNT(*) FROM experiments")
        stats["experiments"] = cursor.fetchone()[0]

        cursor = conn.execute("SELECT COUNT(*) FROM results")
        stats["results"] = cursor.fetchone()[0]

        cursor = conn.execute("SELECT COUNT(*) FROM chunk_embeddings")
        stats["chunks_with_embeddings"] = cursor.fetchone()[0]

        cursor = conn.execute("SELECT COUNT(*) FROM results WHERE is_negative = 1")
        stats["negative_results"] = cursor.fetchone()[0]

        conn.close()
        return stats


# ═══════════════════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════════════════

def main() -> None:
    import argparse

    parser = argparse.ArgumentParser(description="OMEGA Research Memory")
    parser.add_argument("--lab", "-l", default="default", help="Lab name")

    subparsers = parser.add_subparsers(dest="command")

    ingest_parser = subparsers.add_parser("ingest", help="Ingest documents")
    ingest_parser.add_argument("path", help="File or directory to ingest")
    ingest_parser.add_argument(
        "--project", "-p", default="default", help="Project name"
    )

    search_parser = subparsers.add_parser("search", help="Search memory")
    search_parser.add_argument("query", help="Search query")
    search_parser.add_argument("--project", "-p", help="Filter by project")
    search_parser.add_argument(
        "--no-semantic",
        action="store_true",
        help="Use keyword-only search (no embeddings)",
    )

    hyp_parser = subparsers.add_parser("hypothesis", help="Manage hypotheses")
    hyp_parser.add_argument("action", choices=["list", "create", "show"])
    hyp_parser.add_argument("--id", help="Hypothesis ID")
    hyp_parser.add_argument("--claim", help="Claim for new hypothesis")

    subparsers.add_parser("brief", help="Generate weekly brief")
    subparsers.add_parser("stats", help="Show statistics")

    args = parser.parse_args()

    lab_name = getattr(args, "lab", "default")
    memory = ResearchMemory(lab_name)

    if args.command == "ingest":
        count = memory.ingest(args.path, project=args.project)
        print(f"Ingested {count} chunks from {args.path}")

    elif args.command == "search":
        use_semantic = not getattr(args, "no_semantic", False)
        results = memory.query_with_sources(
            args.query,
            project=getattr(args, "project", None),
            use_semantic=use_semantic,
        )
        print(f"\nAnswer:\n{results['answer']}")
        if results["speculation"]:
            print(f"\nSpeculation: {results['speculation_reason']}")

    elif args.command == "hypothesis":
        if args.action == "list":
            hyps = memory.get_hypotheses()
            print(f"\nHypotheses ({len(hyps)} total):\n")
            for h in hyps:
                print(f"  [{h.status}] {h.id}: {h.claim[:60]}...")

        elif args.action == "create" and getattr(args, "claim", None):
            hyp = memory.create_hypothesis(
                claim=args.claim,
                scope="To be defined",
                assumptions=[],
                predictions=[],
            )
            print(f"Created hypothesis: {hyp.id}")

        elif args.action == "show" and getattr(args, "id", None):
            hyps = memory.get_hypotheses()
            for h in hyps:
                if h.id == args.id:
                    print(json.dumps(h.to_dict(), indent=2))
                    break

    elif args.command == "brief":
        brief = memory.generate_weekly_brief()
        print(brief.to_markdown())

    elif args.command == "stats":
        stats = memory.get_stats()
        print(f"\nResearch Memory: {lab_name}")
        print("=" * 40)
        for k, v in stats.items():
            print(f"  {k}: {v}")

    else:
        stats = memory.get_stats()
        print(f"\nResearch Memory: {lab_name}")
        print("=" * 40)
        for k, v in stats.items():
            print(f"  {k}: {v}")


if __name__ == "__main__":
    main()
